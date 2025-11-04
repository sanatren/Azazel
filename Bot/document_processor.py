# Fix the missing client attribute in DocumentProcessor
import os
import tempfile
import logging
from typing import List, Dict, Any, Optional
import docx
from pptx import Presentation
import pdfplumber
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings
from langchain_core.documents import Document
from Bot.vision_processor import VisionProcessor

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process various document types for RAG applications"""

    def __init__(self, api_key=None):
        """Initialize the document processor"""
        # Initialize the text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

        # Use OpenAI embeddings (no heavy dependencies needed)
        self.embeddings = OpenAIEmbeddings(
            model="text-embedding-3-small",  # Cheaper and faster than ada-002
            openai_api_key=api_key or os.getenv("OPENAI_API_KEY")
        )
        print("Using OpenAI embeddings")

        # Dictionary to store vectorstores by session ID
        self.vectorstores = {}

        # Initialize vision processor with provided API key or fallback to env
        if api_key:
            self.vision_processor = VisionProcessor(api_key)
        else:
            self.vision_processor = VisionProcessor(os.getenv("OPENAI_API_KEY"))
    
    def process_file(self, uploaded_file, session_id: str) -> bool:
        """
        Process an uploaded file and store its contents in the vector database

        Args:
            uploaded_file: The uploaded file object (FastAPI UploadFile or Streamlit UploadedFile)
            session_id: The session ID to associate with the document

        Returns:
            bool: True if processing was successful, False otherwise
        """
        try:
            # Check if file is image - handle both FastAPI (content_type) and Streamlit (type)
            file_type = getattr(uploaded_file, 'content_type', getattr(uploaded_file, 'type', ''))
            if file_type.startswith('image/'):
                return self.vision_processor.process_image(uploaded_file, session_id)

            # Create a temporary file to save the uploaded file
            # Handle both FastAPI (read()) and Streamlit (getvalue())
            file_content = uploaded_file.read() if hasattr(uploaded_file, 'read') else uploaded_file.getvalue()
            filename = getattr(uploaded_file, 'filename', getattr(uploaded_file, 'name', 'unknown'))

            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{filename.split('.')[-1]}") as tmp_file:
                tmp_file.write(file_content if isinstance(file_content, bytes) else file_content.read())
                file_path = tmp_file.name
            
            # Extract text based on file type
            file_extension = filename.split('.')[-1].lower()
            
            if file_extension == 'pdf':
                text = self._extract_text_from_pdf(file_path)
            elif file_extension == 'docx':
                text = self._extract_text_from_docx(file_path)
            elif file_extension in ['xlsx', 'xls']:
                text = self._extract_text_from_excel(file_path)
            elif file_extension == 'pptx':
                text = self._extract_text_from_pptx(file_path)
            elif file_extension == 'txt':
                text = self._extract_text_from_txt(file_path)
            else:
                logger.error(f"Unsupported file type: {file_extension}")
                return False
            
            # Clean up the temporary file
            os.unlink(file_path)
            
            if not text:
                logger.warning(f"No text could be extracted from {uploaded_file.name}")
                return False
            
            # Split the text into chunks
            chunks = self.text_splitter.split_text(text)
            
            # Create documents with metadata
            documents = [
                Document(
                    page_content=chunk,
                    metadata={
                        "source": uploaded_file.name,
                        "chunk_id": i,
                        "session_id": session_id
                    }
                )
                for i, chunk in enumerate(chunks)
            ]
            
            # Create or update the vectorstore for this session
            if session_id in self.vectorstores:
                # Add documents to existing vectorstore
                self.vectorstores[session_id].add_documents(documents)
            else:
                # Create a new vectorstore
                self.vectorstores[session_id] = FAISS.from_documents(
                    documents, self.embeddings
                )
            
            return True
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            return False
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from a PDF file"""
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return ""
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from a Word document"""
        try:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting text from Word document: {str(e)}")
            return ""
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from an Excel file"""
        try:
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error extracting text from Excel file: {str(e)}")
            return ""
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from a PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {str(e)}")
            return ""
    
    def _extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from a text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error extracting text from text file: {str(e)}")
            return ""
    
    def query_documents(self, query: str, session_id: str, k: int = 8) -> List[Dict[str, Any]]:
        """
        Query documents for a session
        
        Args:
            query: The query string
            session_id: The session ID
            k: Number of documents to retrieve (increased default from 4 to 8)
            
        Returns:
            List of documents with content and metadata
        """
        try:
            # Check if the vectorstore exists for this session
            if session_id not in self.vectorstores:
                return []
            
            # Get the vectorstore
            vectorstore = self.vectorstores[session_id]
            
            # Query the vectorstore with a higher k to ensure multiple documents are considered
            results = vectorstore.similarity_search_with_score(query, k=k)
            
            # Format the results
            docs = []
            unique_sources = set()  # Track unique document sources
            
            for doc, score in results:
                source = doc.metadata.get("source", "")
                
                # Add to unique sources set
                if source:
                    unique_sources.add(source)
                
                docs.append({
                    "content": doc.page_content,
                    "metadata": doc.metadata,
                    "score": score
                })
            
            # Log the number of unique documents found
            if len(unique_sources) > 0:
                print(f"Retrieved content from {len(unique_sources)} unique documents: {', '.join(unique_sources)}")
            
            return docs
        except Exception as e:
            logger.error(f"Error querying documents: {str(e)}")
            return []
    
    def has_documents(self, session_id: str) -> bool:
        """
        Check if documents exist for a session
        
        Args:
            session_id: The session ID
            
        Returns:
            bool: True if documents exist, False otherwise
        """
        # Simply check if the session_id is in the vectorstores dictionary
        return session_id in self.vectorstores
    
    def clear_documents(self, session_id: str) -> bool:
        """Clear all documents for a session"""
        if session_id in self.vectorstores:
            del self.vectorstores[session_id]
            return True
        return False